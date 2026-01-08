import sys
from pptx import Presentation

def extract_text_from_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text_runs = []
        for i, slide in enumerate(prs.slides):
            text_runs.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
            text_runs.append("\n")
        return "\n".join(text_runs)
    except Exception as e:
        return str(e)

if __name__ == "__main__":
    path = sys.argv[1]
    print(extract_text_from_pptx(path))
